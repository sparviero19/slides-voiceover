"""

"""
import os
import re

import numpy
import pypdf
from pdf2image import convert_from_path
import numpy as np
from pptx import Presentation
from gtts import gTTS
from moviepy.editor import ImageClip, AudioFileClip, concatenate_videoclips

PDF_PATH  = "pdf/CASE2024.pdf"
PPTX_PATH = "pptx/CASE2024.pptx"
TEXT_PATH = "text/CASE2024.md"


def parse_markdown(file_path):
    with open(file_path, 'r') as file:
        content = file.read()
    # Split content into sections based on headers (assuming ## for slide breaks)
    sections = re.split(r'## ', content)
    slides_content = [{'text': section, 'is_code': False} for section in sections[1:]]
    # drop slides
    new_slides_content, slide_index = [], []
    for i, section in enumerate(slides_content):
        if section['text'].find('%-') < 0:
            new_slides_content.append(section)
            slide_index.append(i)
    return new_slides_content, slide_index

def extract_slides_from_pdf(file_path, slide_index=None):
    slides = convert_from_path(file_path)
    if slide_index is None:
        slide_index = list(range(len(slides)))
    new_slides = list()
    for s in slides:
        new_slides.append(np.asarray(s))
    return [new_slides[i] for i in slide_index]


def generate_slide(slide_content):
    prs = Presentation()
    for content in slide_content:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        text_box = slide.shapes.add_textbox(0, 0, width, height)
        text_frame = text_box.text_frame
        text_frame.text = content['text']  # Add text content to slide
    return prs


def convert_text_to_speech(out_filename, text, lang='en'):
    """

    :param out_filename: output filename in mp3 format
    :param text: the string to be converted
    :param lang: the language for tts-ing the string
    :return:
    """
    tts = gTTS(text=text, lang=lang, slow=False)
    tts.save(out_filename)
    return out_filename


def compile_video(slides, tts_audios):
    clips = []
    for slide, audio in zip(slides, tts_audios):
        audio_clip = AudioFileClip(audio)
        slide_clip = ImageClip(slide).set_duration(audio_clip.duration)
        slide_clip = slide_clip.set_audio(audio_clip)
        clips.append(slide_clip)
    final_clip = concatenate_videoclips(clips)
    return final_clip

if __name__ == "__main__":

    scripts, slide_index = parse_markdown(TEXT_PATH)
    audios = list()
    for i, text in enumerate(scripts):
        out = convert_text_to_speech(str(i) + '.mp3', text['text'])
        audios.append(out)

    #slides = Presentation(PPTX_PATH)
    slides = extract_slides_from_pdf(PDF_PATH, slide_index)

    final_video = compile_video(slides, audios)
    output_dir = "outputs"
    final_video.write_videofile(os.path.join(output_dir, "final_video.mp4"), fps=30, codec='libx264')